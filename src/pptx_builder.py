#!/usr/bin/env python3
"""
PPTX Builder
=============
Cross-platform PowerPoint generation using python-pptx.

Inspired by powerpoint-mcp's COM-based approach, but implemented entirely
with python-pptx for Linux/Docker/headless compatibility.

Key capabilities:
  - Template analysis: enumerate layouts and placeholders in .pptx/.potx
  - Placeholder population: fill placeholders with formatted text or images
  - Slide assembly: clone slides from a library into a new proposal deck
  - Speaker notes: attach traceability info (claim IDs, sources)
  - Deck builder: orchestrate full proposal deck from section content

Design decisions borrowed from powerpoint-mcp:
  1. Template-first workflow (analyze before populate)
  2. Placeholder-centric content injection
  3. HTML-like formatting → PPTX runs
  4. Speaker notes for traceability
"""

import copy
import io
import json
import logging
import re
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from xml.etree import ElementTree as ET

logger = logging.getLogger("pptx-builder")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not installed. PPTX builder features unavailable.")


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Foundever brand colors (mirrors frontend/src/lib/brand/tokens.ts)
BRAND_COLORS = {
    "indigo": RGBColor(0x4B, 0x4B, 0xF9) if HAS_PPTX else None,
    "midnight": RGBColor(0x09, 0x09, 0x2D) if HAS_PPTX else None,
    "white": RGBColor(0xFF, 0xFF, 0xFF) if HAS_PPTX else None,
    "light_grey": RGBColor(0xF3, 0xF3, 0xF7) if HAS_PPTX else None,
    "lemon": RGBColor(0xF9, 0xEF, 0x77) if HAS_PPTX else None,
    "coral": RGBColor(0xFF, 0x8D, 0x96) if HAS_PPTX else None,
    "lavender": RGBColor(0xBF, 0xA1, 0xFF) if HAS_PPTX else None,
    "mint": RGBColor(0x8B, 0xF0, 0xBB) if HAS_PPTX else None,
}

# Backend 9-section display names
SECTION_DISPLAY_NAMES = {
    "executive_summary": "Executive Summary",
    "client_understanding": "Understanding Client Needs",
    "solution_overview": "Solution Overview",
    "delivery_model": "Delivery Model",
    "technology": "Technology & Innovation",
    "governance_compliance": "Governance & Compliance",
    "implementation": "Implementation & Transition",
    "team_leadership": "Team & Leadership",
    "proof_points": "Proof Points & Evidence",
}


# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------

@dataclass
class PlaceholderInfo:
    """Information about a single placeholder in a slide layout."""
    idx: int
    name: str
    placeholder_type: str
    left: int
    top: int
    width: int
    height: int

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


@dataclass
class LayoutInfo:
    """Information about a slide layout in a presentation."""
    name: str
    index: int
    placeholders: List[PlaceholderInfo]

    def to_dict(self) -> Dict[str, Any]:
        return {
            "name": self.name,
            "index": self.index,
            "placeholders": [p.to_dict() for p in self.placeholders],
        }


@dataclass
class TemplateAnalysis:
    """Full analysis of a PPTX/POTX template."""
    file_path: str
    slide_width: int
    slide_height: int
    layout_count: int
    layouts: List[LayoutInfo]

    def to_dict(self) -> Dict[str, Any]:
        return {
            "file_path": self.file_path,
            "slide_width": self.slide_width,
            "slide_height": self.slide_height,
            "layout_count": self.layout_count,
            "layouts": [l.to_dict() for l in self.layouts],
        }


@dataclass
class TextSegment:
    """A segment of formatted text to be inserted into a placeholder."""
    text: str
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color: Optional[str] = None  # hex color string
    font_size: Optional[int] = None  # in points


@dataclass
class SlideContent:
    """Content to populate a single slide."""
    section_type: str
    title: str
    body_segments: List[TextSegment] = field(default_factory=list)
    subtitle: Optional[str] = None
    speaker_notes: Optional[str] = None
    image_path: Optional[str] = None
    table_data: Optional[List[List[str]]] = None
    metadata: Dict[str, Any] = field(default_factory=dict)


# ---------------------------------------------------------------------------
# Template Analysis
# ---------------------------------------------------------------------------

def _get_placeholder_type_name(ph) -> str:
    """Get human-readable placeholder type name."""
    try:
        return str(ph.placeholder_format.type).split(".")[-1].split("(")[0]
    except Exception:
        return "UNKNOWN"


def analyze_template(file_path: str) -> TemplateAnalysis:
    """
    Analyze a PPTX/POTX template to discover layouts and placeholders.

    This is the equivalent of powerpoint-mcp's analyze_template tool,
    implemented with python-pptx instead of COM automation.

    Args:
        file_path: Path to .pptx or .potx file

    Returns:
        TemplateAnalysis with all layout and placeholder information.
    """
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not installed")

    prs = Presentation(file_path)
    layouts: List[LayoutInfo] = []

    for i, layout in enumerate(prs.slide_layouts):
        placeholders: List[PlaceholderInfo] = []
        for ph in layout.placeholders:
            placeholders.append(PlaceholderInfo(
                idx=ph.placeholder_format.idx,
                name=ph.name,
                placeholder_type=_get_placeholder_type_name(ph),
                left=ph.left,
                top=ph.top,
                width=ph.width,
                height=ph.height,
            ))

        layouts.append(LayoutInfo(
            name=layout.name,
            index=i,
            placeholders=placeholders,
        ))

    return TemplateAnalysis(
        file_path=file_path,
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
        layout_count=len(layouts),
        layouts=layouts,
    )


# ---------------------------------------------------------------------------
# Text Formatting (HTML-like → PPTX runs)
# ---------------------------------------------------------------------------

# Simplified HTML tag parser (inspired by powerpoint-mcp's process_simple_html)
_TAG_PATTERN = re.compile(r"<(/?)(\w+)(?:\s+[^>]*)?>")


def parse_formatted_text(html_text: str) -> List[TextSegment]:
    """
    Parse simplified HTML-like text into TextSegment list.

    Supported tags:
      <b>bold</b>
      <i>italic</i>
      <u>underline</u>
      <br/> or <br>  (line break)

    Unsupported tags are stripped. Nested tags are handled via state stack.
    """
    segments: List[TextSegment] = []
    bold = False
    italic = False
    underline = False
    pos = 0

    for match in _TAG_PATTERN.finditer(html_text):
        # Add text before this tag
        text_before = html_text[pos:match.start()]
        if text_before:
            segments.append(TextSegment(
                text=text_before,
                bold=bold,
                italic=italic,
                underline=underline,
            ))

        closing = match.group(1) == "/"
        tag = match.group(2).lower()

        if tag == "b":
            bold = not closing
        elif tag == "i":
            italic = not closing
        elif tag == "u":
            underline = not closing
        elif tag == "br":
            segments.append(TextSegment(text="\n"))

        pos = match.end()

    # Add remaining text
    remaining = html_text[pos:]
    if remaining:
        segments.append(TextSegment(
            text=remaining,
            bold=bold,
            italic=italic,
            underline=underline,
        ))

    return segments


def _apply_segments_to_text_frame(text_frame, segments: List[TextSegment]) -> None:
    """Apply TextSegment list to a python-pptx TextFrame."""
    if not segments:
        return

    # Clear existing text
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]

    for i, seg in enumerate(segments):
        if "\n" in seg.text:
            # Split on newlines — each creates a new paragraph
            parts = seg.text.split("\n")
            for j, part in enumerate(parts):
                if j > 0:
                    paragraph = text_frame.add_paragraph()
                if part:
                    run = paragraph.add_run()
                    run.text = part
                    _format_run(run, seg)
        else:
            if i == 0:
                run = paragraph.add_run()
            else:
                run = paragraph.add_run()
            run.text = seg.text
            _format_run(run, seg)


def _format_run(run, seg: TextSegment) -> None:
    """Apply formatting from a TextSegment to a python-pptx Run."""
    run.font.bold = seg.bold
    run.font.italic = seg.italic
    run.font.underline = seg.underline

    if seg.color:
        try:
            color_hex = seg.color.lstrip("#")
            run.font.color.rgb = RGBColor(
                int(color_hex[0:2], 16),
                int(color_hex[2:4], 16),
                int(color_hex[4:6], 16),
            )
        except (ValueError, IndexError):
            pass

    if seg.font_size:
        run.font.size = Pt(seg.font_size)


# ---------------------------------------------------------------------------
# Slide Operations
# ---------------------------------------------------------------------------

def add_slide_from_layout(
    prs,
    layout_index: int = 0,
    position: Optional[int] = None,
) -> Any:
    """
    Add a new slide using a specific layout.

    Args:
        prs: Presentation object
        layout_index: Index of the layout to use
        position: Where to insert (None = end)

    Returns:
        The new slide object
    """
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)

    if position is not None and position < len(prs.slides) - 1:
        # Move the slide to the specified position
        # python-pptx doesn't have a direct move API, so we manipulate XML
        xml_slides = prs.slides._sldIdLst
        slides_list = list(xml_slides)
        new_slide_elem = slides_list[-1]
        xml_slides.remove(new_slide_elem)
        xml_slides.insert(position, new_slide_elem)

    return slide


def populate_placeholder(
    slide,
    placeholder_idx: int,
    content: str,
    formatted: bool = False,
) -> bool:
    """
    Populate a specific placeholder with content.

    Args:
        slide: Slide object
        placeholder_idx: Placeholder index to fill
        content: Text content (plain or HTML-like if formatted=True)
        formatted: Whether to parse HTML-like formatting tags

    Returns:
        True if successful
    """
    try:
        placeholder = slide.placeholders[placeholder_idx]
    except (KeyError, IndexError):
        logger.warning(f"Placeholder {placeholder_idx} not found on slide")
        return False

    if formatted:
        segments = parse_formatted_text(content)
        _apply_segments_to_text_frame(placeholder.text_frame, segments)
    else:
        placeholder.text = content

    return True


def set_speaker_notes(slide, notes_text: str) -> None:
    """
    Set speaker notes on a slide.

    Useful for attaching evidence source IDs, proof tier info,
    and review comments for traceability.
    """
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


def add_table_to_slide(
    slide,
    rows: int,
    cols: int,
    data: List[List[str]],
    left: int = None,
    top: int = None,
    width: int = None,
    height: int = None,
) -> Any:
    """
    Add a table to a slide and populate with data.

    Args:
        slide: Slide object
        rows: Number of rows
        cols: Number of columns
        data: 2D list of cell values
        left/top/width/height: Position in EMU (defaults to centered)

    Returns:
        The table shape
    """
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not installed")

    left = left or Inches(0.5)
    top = top or Inches(1.5)
    width = width or Inches(9.0)
    height = height or Inches(0.3 * rows)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for r_idx, row_data in enumerate(data):
        for c_idx, cell_value in enumerate(row_data):
            if r_idx < rows and c_idx < cols:
                table.cell(r_idx, c_idx).text = str(cell_value)

    return table_shape


# ---------------------------------------------------------------------------
# Slide Cloning (from library source to target deck)
# ---------------------------------------------------------------------------

def clone_slide_from_file(
    target_prs,
    source_path: str,
    source_slide_index: int = 0,
) -> Any:
    """
    Clone a slide from a source PPTX file into the target presentation.

    This is useful for pulling slides from the slide library into a
    proposal deck. Copies shapes, text, and formatting.

    Note: This is a simplified clone that copies shape-by-shape.
    Complex elements (charts, embedded objects) may not transfer perfectly.

    Args:
        target_prs: Target Presentation object
        source_path: Path to source .pptx file
        source_slide_index: Which slide to clone (0-indexed)

    Returns:
        The new slide in the target presentation
    """
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not installed")

    source_prs = Presentation(source_path)

    if source_slide_index >= len(source_prs.slides):
        raise IndexError(
            f"Source has {len(source_prs.slides)} slides, "
            f"requested index {source_slide_index}"
        )

    source_slide = source_prs.slides[source_slide_index]

    # Use the first (blank-ish) layout in target
    blank_layout = None
    for layout in target_prs.slide_layouts:
        if "blank" in layout.name.lower():
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = target_prs.slide_layouts[-1]

    new_slide = target_prs.slides.add_slide(blank_layout)

    # Copy shapes from source to target
    for shape in source_slide.shapes:
        _copy_shape(shape, new_slide)

    # Copy speaker notes if present
    try:
        if source_slide.has_notes_slide:
            notes_text = source_slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                set_speaker_notes(new_slide, notes_text)
    except Exception:
        pass

    return new_slide


def _copy_shape(source_shape, target_slide) -> None:
    """
    Copy a shape from source to target slide.

    Handles text boxes and basic shapes. Tables and charts use
    simplified fallback approaches.
    """
    try:
        # Copy the shape's XML element
        el = copy.deepcopy(source_shape._element)
        target_slide.shapes._spTree.append(el)
    except Exception as e:
        logger.debug(f"Shape copy fallback for {source_shape.shape_type}: {e}")
        # Fallback: if XML copy fails, try to copy text content at least
        if hasattr(source_shape, "text") and source_shape.text.strip():
            try:
                from pptx.util import Inches
                txBox = target_slide.shapes.add_textbox(
                    source_shape.left, source_shape.top,
                    source_shape.width, source_shape.height,
                )
                txBox.text_frame.text = source_shape.text
            except Exception:
                pass


# ---------------------------------------------------------------------------
# Deck Builder — Orchestrate full proposal assembly
# ---------------------------------------------------------------------------

class ProposalDeckBuilder:
    """
    Builds a complete proposal PPTX deck from section content.

    Workflow:
      1. Create or open a template
      2. Add section slides with content
      3. Optionally clone slides from the slide library
      4. Attach speaker notes for traceability
      5. Save the assembled deck

    Usage:
        builder = ProposalDeckBuilder()
        builder.create_from_template("template.pptx")
        builder.add_title_slide("Proposal for {{Client}}", "February 2026")
        builder.add_section_slide(SlideContent(
            section_type="executive_summary",
            title="Executive Summary",
            body_segments=[TextSegment(text="...")],
            speaker_notes="Sources: [T2] Case study...",
        ))
        builder.save("output/proposal.pptx")
    """

    def __init__(self):
        if not HAS_PPTX:
            raise RuntimeError("python-pptx not installed")
        self.prs: Optional[Any] = None
        self.slide_count: int = 0
        self.sections_added: List[str] = []
        self._template_analysis: Optional[TemplateAnalysis] = None

    def create_blank(self) -> None:
        """Create a new blank presentation."""
        self.prs = Presentation()
        self.slide_count = 0
        self.sections_added = []

    def create_from_template(self, template_path: str) -> TemplateAnalysis:
        """
        Create a new presentation from an existing template.

        Also analyzes the template to discover available layouts.
        """
        self.prs = Presentation(template_path)
        self.slide_count = len(self.prs.slides)
        self.sections_added = []
        self._template_analysis = analyze_template(template_path)
        return self._template_analysis

    def open_existing(self, file_path: str) -> int:
        """Open an existing presentation for editing."""
        self.prs = Presentation(file_path)
        self.slide_count = len(self.prs.slides)
        return self.slide_count

    def get_template_analysis(self) -> Optional[TemplateAnalysis]:
        """Return the template analysis (if created from template)."""
        return self._template_analysis

    def find_layout_by_name(self, name_pattern: str) -> Optional[int]:
        """Find a layout index by partial name match."""
        if not self.prs:
            return None
        name_lower = name_pattern.lower()
        for i, layout in enumerate(self.prs.slide_layouts):
            if name_lower in layout.name.lower():
                return i
        return None

    def add_title_slide(
        self,
        title: str,
        subtitle: Optional[str] = None,
        layout_index: int = 0,
    ) -> int:
        """
        Add a title slide to the presentation.

        Returns the slide index (0-based).
        """
        slide = add_slide_from_layout(self.prs, layout_index)

        # Populate title
        if slide.placeholders:
            if 0 in slide.placeholders:
                slide.placeholders[0].text = title
            if subtitle and 1 in slide.placeholders:
                slide.placeholders[1].text = subtitle

        self.slide_count += 1
        return self.slide_count - 1

    def add_section_divider(
        self,
        section_type: str,
        section_number: Optional[int] = None,
    ) -> int:
        """
        Add a section divider slide.

        Uses the section display name from SECTION_DISPLAY_NAMES.
        """
        display_name = SECTION_DISPLAY_NAMES.get(section_type, section_type)
        title = display_name
        if section_number is not None:
            title = f"Section {section_number}: {display_name}"

        # Try to find a "Section" layout, fall back to title layout
        layout_idx = self.find_layout_by_name("section") or 0
        return self.add_title_slide(title, layout_index=layout_idx)

    def add_section_slide(self, content: SlideContent) -> int:
        """
        Add a content slide for a proposal section.

        Uses the best matching layout and populates placeholders
        with the provided content.

        Returns the slide index (0-based).
        """
        # Find appropriate layout
        layout_idx = self.find_layout_by_name("content") or 1
        if layout_idx >= len(self.prs.slide_layouts):
            layout_idx = min(1, len(self.prs.slide_layouts) - 1)

        slide = add_slide_from_layout(self.prs, layout_idx)

        # Populate title placeholder (usually idx 0)
        if 0 in slide.placeholders:
            slide.placeholders[0].text = content.title

        # Populate body placeholder (usually idx 1)
        if content.body_segments and 1 in slide.placeholders:
            _apply_segments_to_text_frame(
                slide.placeholders[1].text_frame,
                content.body_segments,
            )

        # Add table if provided
        if content.table_data and len(content.table_data) > 0:
            rows = len(content.table_data)
            cols = max(len(row) for row in content.table_data)
            add_table_to_slide(slide, rows, cols, content.table_data)

        # Add speaker notes for traceability
        if content.speaker_notes:
            set_speaker_notes(slide, content.speaker_notes)

        self.slide_count += 1
        self.sections_added.append(content.section_type)
        return self.slide_count - 1

    def add_slide_from_library(
        self,
        source_path: str,
        source_slide_index: int = 0,
        speaker_notes_append: Optional[str] = None,
    ) -> int:
        """
        Clone a slide from the slide library into this deck.

        Args:
            source_path: Path to the library slide .pptx
            source_slide_index: Which slide to clone (0-indexed)
            speaker_notes_append: Additional notes to append

        Returns:
            The new slide index
        """
        new_slide = clone_slide_from_file(
            self.prs, source_path, source_slide_index
        )

        if speaker_notes_append:
            try:
                existing_notes = ""
                if new_slide.has_notes_slide:
                    existing_notes = new_slide.notes_slide.notes_text_frame.text
                combined = f"{existing_notes}\n\n{speaker_notes_append}".strip()
                set_speaker_notes(new_slide, combined)
            except Exception:
                set_speaker_notes(new_slide, speaker_notes_append)

        self.slide_count += 1
        return self.slide_count - 1

    def build_full_proposal(
        self,
        sections_content: Dict[str, SlideContent],
        include_dividers: bool = True,
        title: Optional[str] = None,
        subtitle: Optional[str] = None,
    ) -> int:
        """
        Build a complete proposal deck with all sections.

        Args:
            sections_content: Dict mapping backend section ID → content
            include_dividers: Whether to add section divider slides
            title: Proposal title for the cover slide
            subtitle: Subtitle for the cover slide

        Returns:
            Total number of slides
        """
        # Add title slide
        if title:
            self.add_title_slide(title, subtitle)

        # Add sections in order
        section_order = [
            "executive_summary", "client_understanding", "solution_overview",
            "delivery_model", "technology", "governance_compliance",
            "implementation", "team_leadership", "proof_points",
        ]

        for i, section_id in enumerate(section_order, 1):
            if section_id not in sections_content:
                continue

            if include_dividers:
                self.add_section_divider(section_id, section_number=i)

            self.add_section_slide(sections_content[section_id])

        return self.slide_count

    def save(self, output_path: str) -> str:
        """
        Save the presentation to disk.

        Returns the absolute path of the saved file.
        """
        if not self.prs:
            raise RuntimeError("No presentation to save")

        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))

        logger.info(f"Saved proposal deck: {path} ({self.slide_count} slides)")
        return str(path.absolute())

    def save_to_bytes(self) -> bytes:
        """Save the presentation to bytes (for HTTP responses)."""
        if not self.prs:
            raise RuntimeError("No presentation to save")
        buf = io.BytesIO()
        self.prs.save(buf)
        return buf.getvalue()

    def get_build_summary(self) -> Dict[str, Any]:
        """Return summary of the built deck."""
        return {
            "slide_count": self.slide_count,
            "sections_added": self.sections_added,
            "sections_missing": [
                s for s in SECTION_DISPLAY_NAMES
                if s not in self.sections_added
            ],
        }
