#!/usr/bin/env python3
"""
Slide Library Manager
=====================
Indexes, searches, and selects slides from a theme-organized slide library.

Expected library layout:
  slide_library/
    Theme Name/
      slide_001.pptx       # Individual slide or multi-slide deck
      slide_001.txt        # OCR text extract
      slide_002.pptx
      slide_002.txt
    Another Theme/
      ...

Each theme folder may also contain:
  - metadata.json          # Optional: theme-level metadata
  - *.png / *.jpg          # Thumbnail previews

The manager builds an in-memory index at startup (or on demand) and provides
semantic + keyword search across the library, with results mapped to the
3-layer classification taxonomy.
"""

import json
import logging
import re
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger("slide-library")

# ---------------------------------------------------------------------------
# Classification taxonomy constants (mirrors frontend types.ts)
# ---------------------------------------------------------------------------

PRIMARY_LABELS = [
    "executive_summary", "solution_overview", "operational_details",
    "case_study", "compliance_security", "project_plan",
    "pricing", "other", "unclassified",
]

BACKEND_SECTIONS = [
    "executive_summary", "client_understanding", "solution_overview",
    "delivery_model", "technology", "governance_compliance",
    "implementation", "team_leadership", "proof_points",
]

# Keyword signals for classifying slide content to primary labels
LABEL_KEYWORDS: Dict[str, List[str]] = {
    "executive_summary": [
        "executive summary", "overview", "strategic", "partnership",
        "why foundever", "at a glance", "introduction",
    ],
    "solution_overview": [
        "solution", "approach", "capability", "platform", "ecosystem",
        "offering", "service model", "our approach",
    ],
    "operational_details": [
        "process", "workflow", "staffing", "FTE", "headcount", "shift",
        "scorecard", "quality", "site", "facility", "roster", "SOP",
        "workforce", "scheduling", "training", "ramp", "attrition",
        "org chart", "leadership", "account manager",
    ],
    "case_study": [
        "case study", "client example", "success story", "outcome",
        "result", "before and after", "proof point", "testimonial",
    ],
    "compliance_security": [
        "compliance", "security", "certification", "SOC", "PCI",
        "HIPAA", "GDPR", "regulatory", "audit", "risk", "governance",
        "ISO", "NIST",
    ],
    "project_plan": [
        "implementation", "transition", "timeline", "milestone",
        "go-live", "phase", "migration", "ramp plan", "project plan",
    ],
    "pricing": [
        "pricing", "cost", "rate", "commercial", "fee",
        "investment", "budget",
    ],
}

# Section refinement for operational_details fan-out
SECTION_REFINEMENT: Dict[str, Tuple[List[str], str]] = {
    "team_leadership": (
        ["org chart", "leadership", "escalation", "account manager",
         "director", "VP", "management structure"],
        "team_leadership",
    ),
    "technology": (
        ["platform", "tool", "software", "CRM", "telephony", "IVR",
         "API", "AI", "automation", "analytics", "dashboard"],
        "technology",
    ),
    "delivery_model": (
        ["FTE", "staffing", "site", "headcount", "shift", "roster",
         "capacity", "facility", "location", "onshore", "offshore", "nearshore"],
        "delivery_model",
    ),
    "solution_overview": (
        ["process", "workflow", "SOP", "procedure", "methodology"],
        "solution_overview",
    ),
}

# Label → primary backend section mapping
LABEL_TO_PRIMARY_SECTION: Dict[str, str] = {
    "executive_summary": "executive_summary",
    "solution_overview": "solution_overview",
    "operational_details": "delivery_model",
    "case_study": "proof_points",
    "compliance_security": "governance_compliance",
    "project_plan": "implementation",
    "pricing": "executive_summary",
    "other": "executive_summary",
    "unclassified": "executive_summary",
}


# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------

@dataclass
class SlideEntry:
    """A single slide (or deck) in the library."""
    id: str
    theme: str
    file_name: str
    pptx_path: str
    ocr_text: str
    primary_label: str
    backend_section: str
    confidence: float
    keywords_matched: List[str]
    slide_count: int = 1
    has_thumbnail: bool = False
    thumbnail_path: Optional[str] = None
    metadata: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


@dataclass
class ThemeInfo:
    """Summary of a theme folder in the library."""
    name: str
    path: str
    slide_count: int
    label_distribution: Dict[str, int]
    section_distribution: Dict[str, int]
    sample_keywords: List[str]


@dataclass
class SlideSearchResult:
    """A search result from the slide library."""
    slide: SlideEntry
    score: float
    match_reason: str


class SlideLibraryManager:
    """
    Manages a theme-organized slide library.

    Indexes all slides at init, provides keyword and semantic search,
    and maps slides to the 3-layer classification taxonomy.
    """

    def __init__(self, library_path: str):
        self.library_path = Path(library_path)
        self.slides: List[SlideEntry] = []
        self.themes: Dict[str, ThemeInfo] = {}
        self._indexed = False

    def index(self) -> Dict[str, Any]:
        """
        Scan the library directory and build the in-memory index.

        Returns summary stats.
        """
        if not self.library_path.exists():
            raise FileNotFoundError(f"Library path not found: {self.library_path}")

        self.slides = []
        self.themes = {}
        slide_id = 0

        for theme_dir in sorted(self.library_path.iterdir()):
            if not theme_dir.is_dir():
                continue

            theme_name = theme_dir.name
            theme_slides: List[SlideEntry] = []

            # Find all .pptx files in the theme
            pptx_files = sorted(theme_dir.glob("*.pptx"))

            for pptx_path in pptx_files:
                slide_id += 1
                entry_id = f"slide_{slide_id:04d}"

                # Look for matching OCR text file
                ocr_text = ""
                txt_path = pptx_path.with_suffix(".txt")
                if txt_path.exists():
                    ocr_text = txt_path.read_text(encoding="utf-8", errors="replace")

                # Look for thumbnail
                thumbnail_path = None
                has_thumbnail = False
                for ext in [".png", ".jpg", ".jpeg"]:
                    thumb = pptx_path.with_suffix(ext)
                    if thumb.exists():
                        thumbnail_path = str(thumb)
                        has_thumbnail = True
                        break

                # Classify the slide
                label, confidence, matched_kw = self._classify_text(ocr_text, theme_name)
                section = self._resolve_section(label, ocr_text)

                # Get slide count from pptx if python-pptx available
                slide_count = self._count_slides(str(pptx_path))

                # Load optional metadata.json
                meta = {}
                meta_path = pptx_path.with_suffix(".json")
                if meta_path.exists():
                    try:
                        meta = json.loads(meta_path.read_text(encoding="utf-8"))
                    except Exception:
                        pass

                entry = SlideEntry(
                    id=entry_id,
                    theme=theme_name,
                    file_name=pptx_path.name,
                    pptx_path=str(pptx_path),
                    ocr_text=ocr_text,
                    primary_label=label,
                    backend_section=section,
                    confidence=confidence,
                    keywords_matched=matched_kw,
                    slide_count=slide_count,
                    has_thumbnail=has_thumbnail,
                    thumbnail_path=thumbnail_path,
                    metadata=meta,
                )
                theme_slides.append(entry)

            self.slides.extend(theme_slides)

            # Build theme summary
            if theme_slides:
                label_dist: Dict[str, int] = {}
                section_dist: Dict[str, int] = {}
                all_keywords: List[str] = []

                for s in theme_slides:
                    label_dist[s.primary_label] = label_dist.get(s.primary_label, 0) + 1
                    section_dist[s.backend_section] = section_dist.get(s.backend_section, 0) + 1
                    all_keywords.extend(s.keywords_matched[:3])

                self.themes[theme_name] = ThemeInfo(
                    name=theme_name,
                    path=str(theme_dir),
                    slide_count=len(theme_slides),
                    label_distribution=label_dist,
                    section_distribution=section_dist,
                    sample_keywords=list(set(all_keywords))[:10],
                )

        self._indexed = True
        logger.info(
            f"Indexed slide library: {len(self.slides)} slides across "
            f"{len(self.themes)} themes"
        )

        return {
            "total_slides": len(self.slides),
            "total_themes": len(self.themes),
            "themes": {name: info.slide_count for name, info in self.themes.items()},
            "label_distribution": self._global_label_distribution(),
            "section_distribution": self._global_section_distribution(),
        }

    def _classify_text(
        self, text: str, theme_name: str
    ) -> Tuple[str, float, List[str]]:
        """Classify text to a primary label using keyword matching."""
        if not text.strip():
            # Fall back to theme name classification
            text = theme_name

        lower_text = text.lower()
        lower_theme = theme_name.lower()
        combined = f"{lower_theme} {lower_text}"

        scores: Dict[str, Tuple[float, List[str]]] = {}

        for label, keywords in LABEL_KEYWORDS.items():
            matched = [kw for kw in keywords if kw.lower() in combined]
            if matched:
                # Score based on keyword count, normalized by total keywords
                score = len(matched) / len(keywords)
                # Boost if theme name matches
                theme_match = any(kw.lower() in lower_theme for kw in keywords)
                if theme_match:
                    score *= 1.5
                scores[label] = (min(score, 1.0), matched)

        if not scores:
            return "unclassified", 0.0, []

        best_label = max(scores, key=lambda k: scores[k][0])
        confidence, matched_kw = scores[best_label]
        return best_label, round(confidence, 3), matched_kw

    def _resolve_section(self, primary_label: str, text: str) -> str:
        """
        Resolve primary label to a single backend section.
        Handles the operational_details fan-out problem.
        """
        if primary_label != "operational_details":
            return LABEL_TO_PRIMARY_SECTION.get(primary_label, "executive_summary")

        lower_text = text.lower()

        for _section_key, (keywords, section) in SECTION_REFINEMENT.items():
            hits = [kw for kw in keywords if kw.lower() in lower_text]
            if len(hits) >= 2:
                return section

        return "delivery_model"  # Default for operational_details

    def _count_slides(self, pptx_path: str) -> int:
        """Count slides in a PPTX file."""
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            return len(prs.slides)
        except Exception:
            return 1

    def _global_label_distribution(self) -> Dict[str, int]:
        dist: Dict[str, int] = {}
        for s in self.slides:
            dist[s.primary_label] = dist.get(s.primary_label, 0) + 1
        return dist

    def _global_section_distribution(self) -> Dict[str, int]:
        dist: Dict[str, int] = {}
        for s in self.slides:
            dist[s.backend_section] = dist.get(s.backend_section, 0) + 1
        return dist

    # -----------------------------------------------------------------------
    # Search
    # -----------------------------------------------------------------------

    def search(
        self,
        query: str,
        limit: int = 10,
        theme_filter: Optional[str] = None,
        label_filter: Optional[str] = None,
        section_filter: Optional[str] = None,
    ) -> List[SlideSearchResult]:
        """
        Search the slide library by keyword matching against OCR text.

        Returns slides sorted by relevance score (keyword match density).
        """
        self._ensure_indexed()
        query_terms = [t.strip().lower() for t in query.lower().split() if len(t.strip()) >= 2]

        results: List[SlideSearchResult] = []

        for slide in self.slides:
            # Apply filters
            if theme_filter and slide.theme.lower() != theme_filter.lower():
                continue
            if label_filter and slide.primary_label != label_filter:
                continue
            if section_filter and slide.backend_section != section_filter:
                continue

            # Score by keyword match
            search_text = f"{slide.theme} {slide.ocr_text} {slide.file_name}".lower()
            matched_terms = [t for t in query_terms if t in search_text]

            if not matched_terms:
                continue

            score = len(matched_terms) / len(query_terms) if query_terms else 0.0

            # Boost for theme name match
            theme_lower = slide.theme.lower()
            if any(t in theme_lower for t in query_terms):
                score *= 1.3

            # Boost for high-confidence classification
            score *= (0.5 + slide.confidence * 0.5)

            results.append(SlideSearchResult(
                slide=slide,
                score=round(min(score, 1.0), 3),
                match_reason=f"Matched: {', '.join(matched_terms)}",
            ))

        results.sort(key=lambda r: r.score, reverse=True)
        return results[:limit]

    def get_slides_for_section(
        self,
        section: str,
        limit: int = 20,
    ) -> List[SlideEntry]:
        """Get all slides mapped to a specific backend section."""
        self._ensure_indexed()
        matches = [s for s in self.slides if s.backend_section == section]
        matches.sort(key=lambda s: s.confidence, reverse=True)
        return matches[:limit]

    def get_slides_for_label(
        self,
        label: str,
        limit: int = 20,
    ) -> List[SlideEntry]:
        """Get all slides with a specific primary label."""
        self._ensure_indexed()
        matches = [s for s in self.slides if s.primary_label == label]
        matches.sort(key=lambda s: s.confidence, reverse=True)
        return matches[:limit]

    def get_theme_slides(self, theme: str) -> List[SlideEntry]:
        """Get all slides in a specific theme."""
        self._ensure_indexed()
        return [s for s in self.slides if s.theme.lower() == theme.lower()]

    def list_themes(self) -> List[ThemeInfo]:
        """List all themes with summary info."""
        self._ensure_indexed()
        return sorted(self.themes.values(), key=lambda t: t.slide_count, reverse=True)

    def get_library_stats(self) -> Dict[str, Any]:
        """Get overall library statistics."""
        self._ensure_indexed()
        return {
            "library_path": str(self.library_path),
            "total_slides": len(self.slides),
            "total_themes": len(self.themes),
            "themes": {
                name: {
                    "slides": info.slide_count,
                    "labels": info.label_distribution,
                    "sections": info.section_distribution,
                }
                for name, info in self.themes.items()
            },
            "global_label_distribution": self._global_label_distribution(),
            "global_section_distribution": self._global_section_distribution(),
        }

    def select_slides_for_proposal(
        self,
        sections: Optional[List[str]] = None,
        max_per_section: int = 5,
    ) -> Dict[str, List[SlideEntry]]:
        """
        Select best slides for a full proposal, organized by backend section.

        If sections is None, selects for all 9 backend sections.
        Returns dict mapping section → list of best slides.
        """
        self._ensure_indexed()
        target_sections = sections or BACKEND_SECTIONS

        result: Dict[str, List[SlideEntry]] = {}
        for section in target_sections:
            candidates = self.get_slides_for_section(section, limit=max_per_section)
            if candidates:
                result[section] = candidates

        return result

    def _ensure_indexed(self) -> None:
        if not self._indexed:
            self.index()


# ---------------------------------------------------------------------------
# Module-level factory
# ---------------------------------------------------------------------------

_library_manager: Optional[SlideLibraryManager] = None


def get_slide_library(library_path: Optional[str] = None) -> SlideLibraryManager:
    """Get or create the slide library manager singleton."""
    global _library_manager
    if _library_manager is None:
        if library_path is None:
            raise ValueError("library_path required for first initialization")
        _library_manager = SlideLibraryManager(library_path)
    return _library_manager


def set_slide_library_path(library_path: str) -> SlideLibraryManager:
    """Create a new slide library manager for a different path."""
    global _library_manager
    _library_manager = SlideLibraryManager(library_path)
    return _library_manager
