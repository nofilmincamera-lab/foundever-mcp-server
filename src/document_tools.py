#!/usr/bin/env python3
"""
Document Tools Module
=====================
Document manipulation functions for PDF, XLSX, DOCX, PPTX files.
Based on Composio's awesome-claude-skills document-skills.

Auto-saves processed documents to PostgreSQL for training.
"""

import json
import subprocess
import hashlib
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional
from mcp.types import Tool

logger = logging.getLogger("document-tools")

# Database config (matches config.py)
DB_CONFIG = {
    "host": "localhost",
    "port": 5432,
    "database": "bpo_enrichment",
    "user": "bpo_user",
    "password": "bpo_secure_password_2025"
}

# Auto-save toggle (can be disabled if needed)
AUTO_SAVE_ENABLED = True


def save_to_training_db(
    file_path: str,
    file_type: str,
    extracted_text: str,
    extraction_tool: str,
    extracted_tables: Optional[Dict] = None,
    extraction_metadata: Optional[Dict] = None
) -> bool:
    """Save document and extraction to PostgreSQL for training."""
    if not AUTO_SAVE_ENABLED:
        return False

    try:
        import psycopg2
        from psycopg2.extras import Json

        path = Path(file_path)
        if not path.exists():
            return False

        # Read file and compute hash
        with open(file_path, "rb") as f:
            raw_content = f.read()
        file_hash = hashlib.sha256(raw_content).hexdigest()

        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        # Check for duplicates
        cur.execute("SELECT id FROM document_training_data WHERE file_hash = %s", (file_hash,))
        if cur.fetchone():
            logger.info(f"Document already exists (hash: {file_hash[:16]}...)")
            cur.close()
            conn.close()
            return True  # Already saved

        # Insert new record
        cur.execute("""
            INSERT INTO document_training_data
            (filename, file_type, file_size_bytes, file_hash, raw_content,
             extracted_text, extracted_tables, extraction_tool, extraction_metadata, source_path)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
        """, (
            path.name,
            file_type,
            len(raw_content),
            file_hash,
            raw_content,
            extracted_text,
            Json(extracted_tables) if extracted_tables else None,
            extraction_tool,
            Json(extraction_metadata) if extraction_metadata else None,
            str(path.absolute())
        ))

        conn.commit()
        cur.close()
        conn.close()
        logger.info(f"Saved document for training: {path.name} ({file_type})")
        return True

    except Exception as e:
        logger.warning(f"Failed to save to training DB: {e}")
        return False

# ============================================================================
# PDF TOOLS
# ============================================================================

def pdf_extract_text(file_path: str, layout: bool = False) -> str:
    """Extract text from PDF using pdfplumber or pdftotext."""
    extracted_text = ""
    page_count = 0
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text:
                    text_parts.append(f"--- PAGE {i} ---\n{text}")
        extracted_text = "\n\n".join(text_parts) if text_parts else "No text found in PDF"
    except ImportError:
        cmd = ["pdftotext"]
        if layout:
            cmd.append("-layout")
        cmd.extend([file_path, "-"])
        result = subprocess.run(cmd, capture_output=True, text=True)
        extracted_text = result.stdout if result.returncode == 0 else f"Error: {result.stderr}"

    # Auto-save for training
    save_to_training_db(
        file_path=file_path,
        file_type="pdf",
        extracted_text=extracted_text,
        extraction_tool="pdf_extract_text",
        extraction_metadata={"page_count": page_count, "layout": layout}
    )
    return extracted_text


def pdf_extract_tables(file_path: str) -> str:
    """Extract tables from PDF using pdfplumber."""
    try:
        import pdfplumber
        output = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                tables = page.extract_tables()
                if tables:
                    output.append(f"## Page {i}")
                    for j, table in enumerate(tables, 1):
                        output.append(f"### Table {j}")
                        for row in table:
                            output.append(" | ".join(str(cell) if cell else "" for cell in row))
                        output.append("")
        return "\n".join(output) if output else "No tables found in PDF"
    except ImportError:
        return "Error: pdfplumber not installed. Run: pip install pdfplumber"


def pdf_merge(file_paths: List[str], output_path: str) -> str:
    """Merge multiple PDFs into one."""
    try:
        from pypdf import PdfWriter, PdfReader
        writer = PdfWriter()
        for pdf_file in file_paths:
            reader = PdfReader(pdf_file)
            for page in reader.pages:
                writer.add_page(page)
        with open(output_path, "wb") as f:
            writer.write(f)
        return f"Merged {len(file_paths)} PDFs into {output_path}"
    except ImportError:
        return "Error: pypdf not installed. Run: pip install pypdf"


def pdf_split(file_path: str, output_dir: str, pages: Optional[str] = None) -> str:
    """Split PDF into individual pages or ranges."""
    try:
        from pypdf import PdfReader, PdfWriter
        Path(output_dir).mkdir(parents=True, exist_ok=True)
        reader = PdfReader(file_path)

        if pages:
            page_nums = []
            for part in pages.split(","):
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    page_nums.extend(range(start-1, end))
                else:
                    page_nums.append(int(part)-1)
        else:
            page_nums = range(len(reader.pages))

        output_files = []
        for i in page_nums:
            writer = PdfWriter()
            writer.add_page(reader.pages[i])
            out_path = Path(output_dir) / f"page_{i+1}.pdf"
            with open(out_path, "wb") as f:
                writer.write(f)
            output_files.append(str(out_path))

        return f"Split into {len(output_files)} files in {output_dir}"
    except ImportError:
        return "Error: pypdf not installed. Run: pip install pypdf"


def pdf_metadata(file_path: str) -> str:
    """Extract PDF metadata."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        meta = reader.metadata
        info = {
            "Pages": len(reader.pages),
            "Title": meta.title if meta else None,
            "Author": meta.author if meta else None,
            "Subject": meta.subject if meta else None,
            "Creator": meta.creator if meta else None,
        }
        return json.dumps({k: v for k, v in info.items() if v}, indent=2)
    except ImportError:
        return "Error: pypdf not installed"


# ============================================================================
# XLSX TOOLS
# ============================================================================

def xlsx_read(file_path: str, sheet_name: Optional[str] = None) -> str:
    """Read Excel file and return as JSON."""
    try:
        import pandas as pd
        if sheet_name:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            result_json = df.to_json(orient='records', indent=2)
            sheet_names = [sheet_name]
            row_count = len(df)
        else:
            all_sheets = pd.read_excel(file_path, sheet_name=None)
            result = {}
            sheet_names = list(all_sheets.keys())
            row_count = sum(len(df) for df in all_sheets.values())
            for name, df in all_sheets.items():
                result[name] = json.loads(df.to_json(orient='records'))
            result_json = json.dumps(result, indent=2)

        # Auto-save for training
        save_to_training_db(
            file_path=file_path,
            file_type="xlsx",
            extracted_text=result_json,
            extraction_tool="xlsx_read",
            extracted_tables={"sheets": sheet_names},
            extraction_metadata={"sheet_count": len(sheet_names), "total_rows": row_count}
        )
        return result_json
    except ImportError:
        return "Error: pandas not installed. Run: pip install pandas openpyxl"


def xlsx_analyze(file_path: str, sheet_name: Optional[str] = None) -> str:
    """Analyze Excel file - return statistics and structure."""
    try:
        import pandas as pd
        if sheet_name:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            sheets = {sheet_name: df}
        else:
            sheets = pd.read_excel(file_path, sheet_name=None)

        output = []
        for name, df in sheets.items():
            output.append(f"## Sheet: {name}")
            output.append(f"- Rows: {len(df)}")
            output.append(f"- Columns: {len(df.columns)}")
            output.append(f"- Column names: {list(df.columns)}")
            output.append(f"- Data types:\n{df.dtypes.to_string()}")
            output.append(f"- Statistics:\n{df.describe().to_string()}")
            output.append("")

        return "\n".join(output)
    except ImportError:
        return "Error: pandas not installed"


def xlsx_write(data: Dict, output_path: str) -> str:
    """Write data to Excel file."""
    try:
        import pandas as pd

        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            for sheet_name, records in data.items():
                df = pd.DataFrame(records)
                df.to_excel(writer, sheet_name=sheet_name, index=False)

        return f"Written to {output_path}"
    except ImportError:
        return "Error: pandas/openpyxl not installed"


def xlsx_to_csv(file_path: str, output_dir: str) -> str:
    """Convert Excel sheets to CSV files."""
    try:
        import pandas as pd
        Path(output_dir).mkdir(parents=True, exist_ok=True)
        sheets = pd.read_excel(file_path, sheet_name=None)
        output_files = []
        for name, df in sheets.items():
            safe_name = name.replace(" ", "_").replace("/", "_")
            out_path = Path(output_dir) / f"{safe_name}.csv"
            df.to_csv(out_path, index=False)
            output_files.append(str(out_path))
        return f"Converted {len(output_files)} sheets to CSV in {output_dir}"
    except ImportError:
        return "Error: pandas not installed"


# ============================================================================
# DOCX TOOLS
# ============================================================================

def docx_extract_text(file_path: str) -> str:
    """Extract text from Word document."""
    try:
        from docx import Document
        doc = Document(file_path)
        text_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        table_count = len(doc.tables)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(f"[TABLE] {row_text}")

        extracted_text = "\n".join(text_parts)

        # Auto-save for training
        save_to_training_db(
            file_path=file_path,
            file_type="docx",
            extracted_text=extracted_text,
            extraction_tool="docx_extract_text",
            extraction_metadata={"paragraph_count": len(doc.paragraphs), "table_count": table_count}
        )
        return extracted_text
    except ImportError:
        return "Error: python-docx not installed. Run: pip install python-docx"


def docx_to_markdown(file_path: str, track_changes: str = "all") -> str:
    """Convert Word document to markdown using pandoc."""
    cmd = ["pandoc", f"--track-changes={track_changes}", file_path, "-t", "markdown"]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode == 0:
        return result.stdout
    return f"Error: {result.stderr}"


def docx_metadata(file_path: str) -> str:
    """Extract Word document metadata."""
    try:
        from docx import Document
        doc = Document(file_path)
        props = doc.core_properties
        info = {
            "Title": props.title,
            "Author": props.author,
            "Subject": props.subject,
            "Keywords": props.keywords,
            "Created": str(props.created) if props.created else None,
            "Modified": str(props.modified) if props.modified else None,
            "Paragraphs": len(doc.paragraphs),
            "Tables": len(doc.tables),
        }
        return json.dumps({k: v for k, v in info.items() if v}, indent=2)
    except ImportError:
        return "Error: python-docx not installed"


# ============================================================================
# PPTX TOOLS
# ============================================================================

def pptx_extract_text(file_path: str) -> str:
    """Extract text from PowerPoint using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        output = []
        slide_count = len(prs.slides)

        for i, slide in enumerate(prs.slides, 1):
            output.append(f"--- SLIDE {i} ---")
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells if cell.text)
                        if row_text:
                            texts.append(f"[TABLE] {row_text}")
            output.append("\n".join(texts))
            output.append("")

        extracted_text = "\n".join(output)

        # Auto-save for training
        save_to_training_db(
            file_path=file_path,
            file_type="pptx",
            extracted_text=extracted_text,
            extraction_tool="pptx_extract_text",
            extraction_metadata={"slide_count": slide_count}
        )
        return extracted_text
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"


def pptx_to_markdown(file_path: str) -> str:
    """Convert PowerPoint to markdown using markitdown."""
    result = subprocess.run(
        ["python3", "-m", "markitdown", file_path],
        capture_output=True, text=True
    )
    if result.returncode == 0:
        return result.stdout
    return f"Error: {result.stderr}"


def pptx_metadata(file_path: str) -> str:
    """Extract PowerPoint metadata."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        props = prs.core_properties
        info = {
            "Title": props.title,
            "Author": props.author,
            "Subject": props.subject,
            "Slides": len(prs.slides),
            "Created": str(props.created) if props.created else None,
            "Modified": str(props.modified) if props.modified else None,
        }
        return json.dumps({k: v for k, v in info.items() if v}, indent=2)
    except ImportError:
        return "Error: python-pptx not installed"


# ============================================================================
# MCP TOOL DEFINITIONS
# ============================================================================

DOCUMENT_TOOLS = [
    # PDF Tools
    Tool(
        name="pdf_extract_text",
        description="Extract text content from a PDF file. Optionally preserve layout.",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to PDF file"},
                "layout": {"type": "boolean", "description": "Preserve layout (default: false)", "default": False}
            },
            "required": ["file_path"]
        }
    ),
    Tool(
        name="pdf_extract_tables",
        description="Extract tables from a PDF file as structured text.",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to PDF file"}
            },
            "required": ["file_path"]
        }
    ),
    Tool(
        name="pdf_merge",
        description="Merge multiple PDF files into one.",
        inputSchema={
            "type": "object",
            "properties": {
                "file_paths": {"type": "array", "items": {"type": "string"}, "description": "List of PDF file paths to merge"},
                "output_path": {"type": "string", "description": "Output path for merged PDF"}
            },
            "required": ["file_paths", "output_path"]
        }
    ),
    Tool(
        name="pdf_split",
        description="Split a PDF into individual pages or page ranges.",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to PDF file"},
                "output_dir": {"type": "string", "description": "Output directory for split pages"},
                "pages": {"type": "string", "description": "Page ranges like '1-5,7,10-12' (optional, default: all)"}
            },
            "required": ["file_path", "output_dir"]
        }
    ),
    Tool(
        name="pdf_metadata",
        description="Extract metadata from a PDF file (title, author, pages, etc).",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to PDF file"}
            },
            "required": ["file_path"]
        }
    ),

    # XLSX Tools
    Tool(
        name="xlsx_read",
        description="Read Excel file and return data as JSON. Can read specific sheet or all sheets.",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to Excel file"},
                "sheet_name": {"type": "string", "description": "Specific sheet name (optional, default: all sheets)"}
            },
            "required": ["file_path"]
        }
    ),
    Tool(
        name="xlsx_analyze",
        description="Analyze Excel file - get statistics, structure, column types, and summary stats.",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to Excel file"},
                "sheet_name": {"type": "string", "description": "Specific sheet name (optional)"}
            },
            "required": ["file_path"]
        }
    ),
    Tool(
        name="xlsx_write",
        description="Write data to Excel file. Data should be a dict with sheet names as keys and arrays of records as values.",
        inputSchema={
            "type": "object",
            "properties": {
                "data": {"type": "object", "description": "Dict with sheet names as keys, arrays of records as values"},
                "output_path": {"type": "string", "description": "Output path for Excel file"}
            },
            "required": ["data", "output_path"]
        }
    ),
    Tool(
        name="xlsx_to_csv",
        description="Convert Excel sheets to CSV files.",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to Excel file"},
                "output_dir": {"type": "string", "description": "Output directory for CSV files"}
            },
            "required": ["file_path", "output_dir"]
        }
    ),

    # DOCX Tools
    Tool(
        name="docx_extract_text",
        description="Extract text content from a Word document including tables.",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to Word document"}
            },
            "required": ["file_path"]
        }
    ),
    Tool(
        name="docx_to_markdown",
        description="Convert Word document to markdown using pandoc. Can preserve tracked changes.",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to Word document"},
                "track_changes": {"type": "string", "enum": ["accept", "reject", "all"], "default": "all"}
            },
            "required": ["file_path"]
        }
    ),
    Tool(
        name="docx_metadata",
        description="Extract metadata from Word document (author, title, dates, etc).",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to Word document"}
            },
            "required": ["file_path"]
        }
    ),

    # PPTX Tools
    Tool(
        name="pptx_extract_text",
        description="Extract text content from PowerPoint slides including tables.",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to PowerPoint file"}
            },
            "required": ["file_path"]
        }
    ),
    Tool(
        name="pptx_to_markdown",
        description="Convert PowerPoint to markdown using markitdown.",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to PowerPoint file"}
            },
            "required": ["file_path"]
        }
    ),
    Tool(
        name="pptx_metadata",
        description="Extract metadata from PowerPoint file.",
        inputSchema={
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Path to PowerPoint file"}
            },
            "required": ["file_path"]
        }
    ),
]


async def handle_document_tool(name: str, arguments: Dict[str, Any]) -> str:
    """Handle document tool calls."""
    # PDF Tools
    if name == "pdf_extract_text":
        return pdf_extract_text(arguments["file_path"], arguments.get("layout", False))
    elif name == "pdf_extract_tables":
        return pdf_extract_tables(arguments["file_path"])
    elif name == "pdf_merge":
        return pdf_merge(arguments["file_paths"], arguments["output_path"])
    elif name == "pdf_split":
        return pdf_split(arguments["file_path"], arguments["output_dir"], arguments.get("pages"))
    elif name == "pdf_metadata":
        return pdf_metadata(arguments["file_path"])

    # XLSX Tools
    elif name == "xlsx_read":
        return xlsx_read(arguments["file_path"], arguments.get("sheet_name"))
    elif name == "xlsx_analyze":
        return xlsx_analyze(arguments["file_path"], arguments.get("sheet_name"))
    elif name == "xlsx_write":
        return xlsx_write(arguments["data"], arguments["output_path"])
    elif name == "xlsx_to_csv":
        return xlsx_to_csv(arguments["file_path"], arguments["output_dir"])

    # DOCX Tools
    elif name == "docx_extract_text":
        return docx_extract_text(arguments["file_path"])
    elif name == "docx_to_markdown":
        return docx_to_markdown(arguments["file_path"], arguments.get("track_changes", "all"))
    elif name == "docx_metadata":
        return docx_metadata(arguments["file_path"])

    # PPTX Tools
    elif name == "pptx_extract_text":
        return pptx_extract_text(arguments["file_path"])
    elif name == "pptx_to_markdown":
        return pptx_to_markdown(arguments["file_path"])
    elif name == "pptx_metadata":
        return pptx_metadata(arguments["file_path"])

    return None  # Not a document tool
